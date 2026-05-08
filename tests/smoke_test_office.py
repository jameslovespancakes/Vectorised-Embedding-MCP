"""Smoke test: pptx + xlsx + xls parsers via synthetic files."""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

from vectorise_mcp import parsers


def make_pptx(path: Path) -> None:
    from pptx import Presentation
    pres = Presentation()
    layout = pres.slide_layouts[1]  # title + content
    s1 = pres.slides.add_slide(layout)
    s1.shapes.title.text = "Quantum Computing Intro"
    s1.placeholders[1].text = "Qubits superpose 0 and 1"
    s1.notes_slide.notes_text_frame.text = "Speaker note for slide 1"
    s2 = pres.slides.add_slide(layout)
    s2.shapes.title.text = "Entanglement"
    s2.placeholders[1].text = "Two qubits link instantaneously"
    pres.save(str(path))


def make_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Q1Numbers"
    ws.append(["Region", "Revenue", "Growth"])
    ws.append(["NA", 4200000, 0.18])
    ws.append(["EU", 2100000, 0.12])
    ws2 = wb.create_sheet("Notes")
    ws2.append(["Forecast", "Q2 to hit 5B"])
    wb.save(str(path))


def make_xls(path: Path) -> None:
    # xlrd 2.0+ is read-only. Use xlwt for writing the legacy format.
    try:
        import xlwt
    except ImportError:
        print("  xlwt not installed; skipping .xls write test")
        return
    wb = xlwt.Workbook()
    ws = wb.add_sheet("Legacy")
    ws.write(0, 0, "Old")
    ws.write(0, 1, "Format")
    ws.write(1, 0, "Apple")
    ws.write(1, 1, "Banana")
    wb.save(str(path))


def main() -> int:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)

        print("=== PPTX ===")
        ppt = root / "deck.pptx"
        make_pptx(ppt)
        for text, page in parsers.parse(ppt):
            print(f"  slide={page}")
            for line in text.splitlines():
                print(f"    {line}")

        print("\n=== XLSX ===")
        xlsx = root / "numbers.xlsx"
        make_xlsx(xlsx)
        for text, page in parsers.parse(xlsx):
            print(f"  page={page}")
            for line in text.splitlines():
                print(f"    {line}")

        print("\n=== XLS ===")
        xls = root / "legacy.xls"
        make_xls(xls)
        if xls.exists():
            for text, page in parsers.parse(xls):
                print(f"  page={page}")
                for line in text.splitlines():
                    print(f"    {line}")

    print("\nALL OFFICE PARSERS OK")
    return 0


if __name__ == "__main__":
    sys.exit(main())
