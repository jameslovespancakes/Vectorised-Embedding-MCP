"""PowerPoint parser via python-pptx. Yields (text, slide_number) per slide.

Captures all text-bearing shapes (titles, content placeholders, text boxes,
tables) plus speaker notes. Uses the slide number as the `page` field so
metadata filters and search results match users' mental model of "slide N".
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Iterator

logger = logging.getLogger(__name__)


def _shape_text(shape) -> str:
    parts: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)
    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def parse(path: Path) -> Iterator[tuple[str, int]]:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — skipping %s", path)
        return

    pres = Presentation(str(path))
    for slide_num, slide in enumerate(pres.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            try:
                t = _shape_text(shape)
                if t:
                    parts.append(t)
            except Exception:
                logger.debug("shape extract failed", exc_info=True)

        # Speaker notes
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker notes]\n{notes}")
        except Exception:
            logger.debug("notes extract failed", exc_info=True)

        text = "\n\n".join(parts).strip()
        if text:
            yield text, slide_num
